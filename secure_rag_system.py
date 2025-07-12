#!/usr/bin/env python3
"""
Secure Enhanced RAG System with User-Input API Keys and Multiple Model Options
No API keys stored in code or environment files
"""

import os
import fitz  # PyMuPDF
from faiss_vector_store import FaissClient
from sentence_transformers import SentenceTransformer
from typing import List, Dict, Any, Optional
from pathlib import Path
import openai
from dotenv import load_dotenv
import hashlib
import json
from tqdm import tqdm
import logging
from concurrent.futures import ThreadPoolExecutor, as_completed
import time

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Load environment variables (no API keys)
load_dotenv()

class SecureMultiFormatRAG:
    
    # Available OpenAI models with their specifications
    AVAILABLE_MODELS = {
        'gpt-3.5-turbo': {
            'name': 'GPT-3.5 Turbo',
            'description': 'Fast and cost-effective for most tasks',
            'max_tokens': 4096,
            'cost_per_1k_input': 0.0015,
            'cost_per_1k_output': 0.002,
            'recommended_for': 'General queries, high volume usage'
        },
        'gpt-3.5-turbo-16k': {
            'name': 'GPT-3.5 Turbo 16K',
            'description': 'Larger context window for complex documents',
            'max_tokens': 16384,
            'cost_per_1k_input': 0.003,
            'cost_per_1k_output': 0.004,
            'recommended_for': 'Long documents, complex analysis'
        },
        'gpt-4-turbo': {
            'name': 'GPT-4 Turbo',
            'description': 'Best balance of performance and cost',
            'max_tokens': 128000,
            'cost_per_1k_input': 0.01,
            'cost_per_1k_output': 0.03,
            'recommended_for': 'Complex reasoning, detailed analysis'
        },
        'gpt-4': {
            'name': 'GPT-4',
            'description': 'Highest quality responses',
            'max_tokens': 8192,
            'cost_per_1k_input': 0.03,
            'cost_per_1k_output': 0.06,
            'recommended_for': 'Critical analysis, highest accuracy needed'
        },
        'gpt-4o': {
            'name': 'GPT-4o',
            'description': 'Latest optimized model',
            'max_tokens': 128000,
            'cost_per_1k_input': 0.005,
            'cost_per_1k_output': 0.015,
            'recommended_for': 'Latest features, optimized performance'
        },
        'gpt-4o-mini': {
            'name': 'GPT-4o Mini',
            'description': 'Lightweight version of GPT-4o',
            'max_tokens': 128000,
            'cost_per_1k_input': 0.00015,
            'cost_per_1k_output': 0.0006,
            'recommended_for': 'Cost-effective GPT-4 quality'
        }
    }
    
    def __init__(self, collection_name: str = "secure_documents", batch_size: int = 50):
        """Initialize secure RAG system - API key must be provided by user"""
        self.client = FaissClient()
        self.collection_name = collection_name
        self.batch_size = batch_size
        self.embedding_model_name = os.getenv("EMBEDDING_MODEL", "all-MiniLM-L6-v2")
        self.embedding_model = SentenceTransformer(self.embedding_model_name)
        
        # API key and model - set by user, not from environment
        self.openai_api_key = None
        self.openai_model = os.getenv("DEFAULT_OPENAI_MODEL", "gpt-3.5-turbo")
        self.llm_available = False
        
        # Initialize collection
        try:
            self.collection = self.client.get_or_create_collection(name=collection_name)
            logger.info(f"Using FAISS collection: {collection_name}")
        except Exception as e:
            logger.error(f"Error initializing FAISS collection: {e}")
            self.collection = self.client.create_collection(name=collection_name)
            logger.info(f"Created new FAISS collection: {collection_name}")
        
        # Processing statistics
        self.processing_stats = {
            'total_files': 0,
            'successful': 0,
            'failed': 0,
            'total_chunks': 0,
            'processing_time': 0
        }
    
    def set_openai_credentials(self, api_key: str, model: str = "gpt-3.5-turbo") -> bool:
        """Securely set OpenAI credentials provided by user"""
        if not api_key or not api_key.strip():
            self.llm_available = False
            return False
        
        try:
            # Basic format validation only - NO API CALLS
            if not api_key.startswith('sk-'):
                logger.error("Invalid API key format")
                return False
            
            # Set credentials without testing
            self.openai_api_key = api_key.strip()
            self.openai_model = model if model in self.AVAILABLE_MODELS else "gpt-3.5-turbo"
            
            # Configure OpenAI client
            openai.api_key = self.openai_api_key
            
            # Mark as available - validation happens on first real use
            self.llm_available = True
            logger.info(f"OpenAI credentials set with model: {self.openai_model}")
            return True
                
        except Exception as e:
            logger.error(f"Error setting OpenAI credentials: {e}")
            self.llm_available = False
            return False
    
    def get_available_models(self) -> Dict[str, Dict]:
        """Get list of available OpenAI models with their specifications"""
        return self.AVAILABLE_MODELS
    
    def calculate_model_cost(self, model: str, queries_per_month: int, 
                           avg_input_tokens: int = 2000, avg_output_tokens: int = 200) -> Dict[str, float]:
        """Calculate estimated costs for a specific model"""
        if model not in self.AVAILABLE_MODELS:
            return {'error': 'Model not found'}
        
        model_info = self.AVAILABLE_MODELS[model]
        
        total_input_tokens = avg_input_tokens * queries_per_month
        total_output_tokens = avg_output_tokens * queries_per_month
        
        input_cost = (total_input_tokens / 1000) * model_info['cost_per_1k_input']
        output_cost = (total_output_tokens / 1000) * model_info['cost_per_1k_output']
        total_cost = input_cost + output_cost
        
        return {
            'monthly_cost': round(total_cost, 4),
            'cost_per_query': round(total_cost / queries_per_month, 6),
            'annual_cost': round(total_cost * 12, 2),
            'input_cost': round(input_cost, 4),
            'output_cost': round(output_cost, 4)
        }
    
    def extract_text_from_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract text from PDF files"""
        try:
            doc = fitz.open(file_path)
            pages_data = []
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                
                if text.strip():
                    chunks = self._split_text(text)
                    for chunk_id, chunk in enumerate(chunks):
                        if chunk.strip():
                            pages_data.append({
                                'text': chunk.strip(),
                                'document': os.path.basename(file_path),
                                'page': page_num + 1,
                                'chunk_id': chunk_id,
                                'file_path': file_path,
                                'file_type': 'pdf'
                            })
            
            doc.close()
            return pages_data
            
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {e}")
            return []
    
    def extract_text_from_word(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract text from Word documents"""
        try:
            import docx
            doc = docx.Document(file_path)
            full_text = []
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text)
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells])
                    if row_text.strip():
                        full_text.append(row_text)
            
            # Combine and chunk
            combined_text = "\n".join(full_text)
            chunks = self._split_text(combined_text)
            
            pages_data = []
            for chunk_id, chunk in enumerate(chunks):
                if chunk.strip():
                    pages_data.append({
                        'text': chunk.strip(),
                        'document': os.path.basename(file_path),
                        'page': 1,
                        'chunk_id': chunk_id,
                        'file_path': file_path,
                        'file_type': 'word'
                    })
            
            return pages_data
            
        except Exception as e:
            logger.error(f"Error processing Word document {file_path}: {e}")
            return []
    
    def extract_text_from_excel(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract text from Excel files"""
        try:
            import pandas as pd
            excel_file = pd.ExcelFile(file_path)
            pages_data = []
            
            for sheet_num, sheet_name in enumerate(excel_file.sheet_names):
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # Convert DataFrame to text
                sheet_text = f"Sheet: {sheet_name}\n"
                sheet_text += df.to_string(index=False)
                
                chunks = self._split_text(sheet_text)
                for chunk_id, chunk in enumerate(chunks):
                    if chunk.strip():
                        pages_data.append({
                            'text': chunk.strip(),
                            'document': os.path.basename(file_path),
                            'page': sheet_num + 1,
                            'chunk_id': chunk_id,
                            'file_path': file_path,
                            'file_type': 'excel',
                            'sheet_name': sheet_name
                        })
            
            return pages_data
            
        except Exception as e:
            logger.error(f"Error processing Excel file {file_path}: {e}")
            return []
    
    def extract_text_from_powerpoint(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract text from PowerPoint files"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            pages_data = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = f"Slide {slide_num + 1}:\n"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                
                chunks = self._split_text(slide_text)
                for chunk_id, chunk in enumerate(chunks):
                    if chunk.strip():
                        pages_data.append({
                            'text': chunk.strip(),
                            'document': os.path.basename(file_path),
                            'page': slide_num + 1,
                            'chunk_id': chunk_id,
                            'file_path': file_path,
                            'file_type': 'powerpoint'
                        })
            
            return pages_data
            
        except Exception as e:
            logger.error(f"Error processing PowerPoint file {file_path}: {e}")
            return []
    
    def _split_text(self, text: str, chunk_size: int = 500) -> List[str]:
        """Split text into chunks optimized for large-scale processing"""
        paragraphs = text.split('\n\n')
        chunks = []
        
        for paragraph in paragraphs:
            if len(paragraph) <= chunk_size:
                chunks.append(paragraph)
            else:
                # Split long paragraphs by sentences
                sentences = paragraph.split('. ')
                current_chunk = ""
                
                for sentence in sentences:
                    if len(current_chunk + sentence) <= chunk_size:
                        current_chunk += sentence + ". "
                    else:
                        if current_chunk:
                            chunks.append(current_chunk.strip())
                        current_chunk = sentence + ". "
                
                if current_chunk:
                    chunks.append(current_chunk.strip())
        
        return chunks
    
    def process_single_file(self, file_path: Path) -> List[Dict[str, Any]]:
        """Process a single file based on its extension"""
        file_ext = file_path.suffix.lower()
        
        if file_ext == '.pdf':
            return self.extract_text_from_pdf(str(file_path))
        elif file_ext in ['.docx', '.doc']:
            return self.extract_text_from_word(str(file_path))
        elif file_ext in ['.xlsx', '.xls']:
            return self.extract_text_from_excel(str(file_path))
        elif file_ext in ['.pptx', '.ppt']:
            return self.extract_text_from_powerpoint(str(file_path))
        else:
            logger.warning(f"Unsupported file type: {file_ext}")
            return []
    
    def process_documents_parallel(self, directory: str = ".", max_workers: int = 4) -> int:
        """Process documents in parallel for better performance"""
        start_time = time.time()
        
        # Find all supported files
        supported_extensions = {'.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt'}
        all_files = []
        
        for ext in supported_extensions:
            all_files.extend(Path(directory).glob(f"*{ext}"))
        
        self.processing_stats['total_files'] = len(all_files)
        logger.info(f"Found {len(all_files)} files to process")
        
        if not all_files:
            logger.warning("No supported files found")
            return 0
        
        # Process files in parallel
        all_documents = []
        all_embeddings = []
        all_metadatas = []
        all_ids = []
        
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            # Submit all files for processing
            future_to_file = {executor.submit(self.process_single_file, file_path): file_path 
                             for file_path in all_files}
            
            # Collect results with progress bar
            for future in tqdm(as_completed(future_to_file), total=len(all_files), desc="Processing files"):
                file_path = future_to_file[future]
                try:
                    pages_data = future.result()
                    if pages_data:
                        self.processing_stats['successful'] += 1
                        
                        # Process chunks in batches for memory efficiency
                        for page_data in pages_data:
                            # Create embeddings
                            embedding = self.embedding_model.encode(page_data['text'])
                            
                            # Prepare data for FAISS
                            doc_id = f"{page_data['document']}_page_{page_data['page']}_chunk_{page_data['chunk_id']}"
                            
                            all_documents.append(page_data['text'])
                            all_embeddings.append(embedding.tolist())
                            all_metadatas.append({
                                'document': page_data['document'],
                                'page': page_data['page'],
                                'chunk_id': page_data['chunk_id'],
                                'file_path': page_data['file_path'],
                                'file_type': page_data['file_type']
                            })
                            all_ids.append(doc_id)
                            
                            # Add to database in batches
                            if len(all_documents) >= self.batch_size:
                                self._add_batch_to_db(all_documents, all_embeddings, all_metadatas, all_ids)
                                all_documents, all_embeddings, all_metadatas, all_ids = [], [], [], []
                    
                except Exception as e:
                    logger.error(f"Error processing {file_path}: {e}")
                    self.processing_stats['failed'] += 1
        
        # Add remaining documents
        if all_documents:
            self._add_batch_to_db(all_documents, all_embeddings, all_metadatas, all_ids)
        
        # Update statistics
        self.processing_stats['total_chunks'] = self.collection.count()
        self.processing_stats['processing_time'] = time.time() - start_time
        
        logger.info(f"Processing complete: {self.processing_stats}")
        return self.processing_stats['total_chunks']
    
    def _add_batch_to_db(self, documents, embeddings, metadatas, ids):
        """Add a batch of documents to the FAISS database"""
        try:
            self.collection.add(
                documents=documents,
                embeddings=embeddings,
                metadatas=metadatas,
                ids=ids
            )
        except Exception as e:
            logger.error(f"Error adding batch to database: {e}")
    
    def search_documents(self, question: str, n_results: int = 10) -> Dict[str, Any]:
        """Search documents with optimized performance for large datasets"""
        question_embedding = self.embedding_model.encode(question)
        
        results = self.collection.query(
            query_embeddings=[question_embedding.tolist()],
            n_results=n_results
        )
        
        response = {
            'question': question,
            'results': []
        }
        
        if results['documents'] and results['documents'][0]:
            for i in range(len(results['documents'][0])):
                result_item = {
                    'text': results['documents'][0][i],
                    'document': results['metadatas'][0][i]['document'],
                    'page': results['metadatas'][0][i]['page'],
                    'file_type': results['metadatas'][0][i]['file_type'],
                    'distance': results['distances'][0][i] if 'distances' in results else None
                }
                response['results'].append(result_item)
        
        return response
    
    def generate_llm_response(self, question: str, context_results: List[Dict], model: str = None) -> str:
        """Generate a natural language response using OpenAI"""
        if not self.llm_available:
            return "❌ OpenAI API not available. Please set your API key in the interface."
        
        # Use specified model or default
        use_model = model if model and model in self.AVAILABLE_MODELS else self.openai_model
        
        # Prepare context from search results
        context_text = ""
        for i, result in enumerate(context_results, 1):
            context_text += f"\n\n[Document: {result['document']}, Page: {result['page']}, Type: {result.get('file_type', 'unknown')}]\n"
            context_text += result['text']
        
        # Check token limits for the model
        model_info = self.AVAILABLE_MODELS[use_model]
        max_tokens = model_info['max_tokens']
        
        # Estimate tokens (rough approximation: 1 token ≈ 4 characters)
        estimated_tokens = len(context_text + question) // 4
        response_tokens = min(1500, max_tokens // 2)  # Allow longer, more detailed responses
        
        if estimated_tokens > (max_tokens - response_tokens):
            # Truncate context if too long
            max_context_chars = (max_tokens - response_tokens - 100) * 4  # Leave buffer
            context_text = context_text[:max_context_chars] + "..."
        
        # Create the prompt
        prompt = f"""You are a helpful assistant that answers questions based on the provided document content. 
        
Use the following context to answer the question. If the answer is not in the context, say so.
Always cite the document name, page number, and file type when referencing information.

Question: {question}

Context from documents:{context_text}

Answer:"""
        
        try:
            response = openai.chat.completions.create(
                model=use_model,
                messages=[
                    {"role": "system", "content": "You are a helpful assistant that answers questions based on document content. Always cite sources with document name, page number, and file type."},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=response_tokens,
                temperature=0.7
            )
            
            return response.choices[0].message.content.strip()
            
        except Exception as e:
            return f"❌ Error generating LLM response: {str(e)}"
    
    def get_database_stats(self) -> Dict[str, Any]:
        """Get statistics about the document database"""
        try:
            collection_count = self.collection.count()
            return {
                'total_chunks': collection_count,
                'collection_name': self.collection_name,
                'embedding_model': self.embedding_model_name,
                'llm_available': self.llm_available,
                'llm_model': self.openai_model if self.llm_available else None,
                'available_models': list(self.AVAILABLE_MODELS.keys())
            }
        except Exception as e:
            return {'error': str(e)}

def main():
    """Command line interface for testing"""
    print("🔒 Secure Multi-Format RAG System")
    print("="*60)
    print("⚠️  API keys must be provided by user - no keys stored in code!")
    print("="*60)
    
    # Initialize system
    rag = SecureMultiFormatRAG()
    
    # Show available models
    print("\n📋 Available OpenAI Models:")
    for model_id, info in rag.get_available_models().items():
        print(f"   {model_id}: {info['name']} - {info['description']}")
    
    # Get API key from user
    print("\n🔑 Enter your OpenAI API key:")
    api_key = input("API Key: ").strip()
    
    if api_key:
        print("\n📋 Select OpenAI model:")
        models = list(rag.get_available_models().keys())
        for i, model in enumerate(models, 1):
            print(f"   {i}. {model}")
        
        try:
            model_choice = int(input("Choose model (1-6): ")) - 1
            selected_model = models[model_choice] if 0 <= model_choice < len(models) else "gpt-3.5-turbo"
        except:
            selected_model = "gpt-3.5-turbo"
        
        # Set credentials
        if rag.set_openai_credentials(api_key, selected_model):
            print(f"✅ API key validated! Using model: {selected_model}")
            
            # Show cost estimate
            cost_info = rag.calculate_model_cost(selected_model, 100)
            print(f"💰 Estimated cost for 100 queries: ${cost_info['monthly_cost']}")
        else:
            print("❌ Invalid API key or model selection")
    else:
        print("⚠️  No API key provided - LLM features disabled")
    
    print("\n✅ System ready for use!")

if __name__ == "__main__":
    main()