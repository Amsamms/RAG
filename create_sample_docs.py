#!/usr/bin/env python3
"""
Create sample documents for testing the multi-format RAG system
"""

import pandas as pd
from docx import Document
from pptx import Presentation
from pptx.util import Inches

def create_sample_word_doc():
    """Create a sample Word document"""
    doc = Document()
    
    # Title
    title = doc.add_heading('Chemical Process Safety Manual', 0)
    
    # Introduction
    doc.add_heading('Introduction', level=1)
    intro = doc.add_paragraph(
        'This manual covers essential chemical processes and safety procedures for industrial operations. '
        'It includes information about various chemical compounds, their properties, and handling procedures.'
    )
    
    # Chapter 1: Chemical Compounds
    doc.add_heading('Chapter 1: Chemical Compounds', level=1)
    
    doc.add_heading('Ethanolamines', level=2)
    ethanolamine_para = doc.add_paragraph(
        'Ethanolamines are important industrial chemicals used in gas processing and chemical synthesis. '
        'They include monoethanolamine (MEA), diethanolamine (DEA), and triethanolamine (TEA). '
        'These compounds are commonly used for acid gas removal in natural gas processing plants.'
    )
    
    doc.add_heading('Safety Procedures', level=2)
    safety_para = doc.add_paragraph(
        'When handling ethanolamines, proper personal protective equipment must be worn. '
        'Ensure adequate ventilation and follow all safety protocols outlined in the MSDS sheets.'
    )
    
    # Chapter 2: Process Operations
    doc.add_heading('Chapter 2: Process Operations', level=1)
    
    doc.add_heading('Gas Processing', level=2)
    gas_para = doc.add_paragraph(
        'Gas processing involves several steps including acid gas removal, dehydration, and fractionation. '
        'The Merox process is commonly used for sulfur compound removal from hydrocarbon streams.'
    )
    
    # Table
    table = doc.add_table(rows=1, cols=3)
    table.style = 'Table Grid'
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = 'Chemical'
    hdr_cells[1].text = 'Formula'
    hdr_cells[2].text = 'Use'
    
    # Add rows
    chemicals = [
        ('Monoethanolamine', 'C2H7NO', 'Acid gas removal'),
        ('Diethanolamine', 'C4H11NO2', 'Gas treating'),
        ('Methyldiethanolamine', 'C5H13NO2', 'Selective H2S removal')
    ]
    
    for chem, formula, use in chemicals:
        row_cells = table.add_row().cells
        row_cells[0].text = chem
        row_cells[1].text = formula
        row_cells[2].text = use
    
    doc.save('sample_chemical_manual.docx')
    print("✅ Created sample_chemical_manual.docx")

def create_sample_excel_file():
    """Create a sample Excel file with multiple sheets"""
    
    # Chemical properties sheet
    chemicals_data = {
        'Chemical Name': [
            'Monoethanolamine (MEA)',
            'Diethanolamine (DEA)', 
            'Triethanolamine (TEA)',
            'Methyldiethanolamine (MDEA)',
            'Diisopropanolamine (DIPA)'
        ],
        'Formula': ['C2H7NO', 'C4H11NO2', 'C6H15NO3', 'C5H13NO2', 'C6H15NO2'],
        'Molecular Weight': [61.08, 105.14, 149.19, 119.16, 133.19],
        'Boiling Point (°C)': [170, 268, 335, 247, 249],
        'Density (g/mL)': [1.012, 1.097, 1.124, 1.043, 0.989],
        'Primary Use': [
            'Acid gas removal',
            'Gas treating, sulfur removal', 
            'Gas treating, corrosion inhibitor',
            'Selective H2S removal',
            'Gas treating'
        ]
    }
    
    # Process conditions sheet
    process_data = {
        'Process': ['Merox', 'Amine Treating', 'Dehydration', 'Fractionation'],
        'Temperature Range (°C)': ['20-60', '40-120', '150-200', '80-150'],
        'Pressure Range (bar)': ['1-10', '5-50', '10-30', '2-20'],
        'Key Chemicals': [
            'Caustic, Oxygen',
            'MEA, DEA, MDEA',
            'Triethylene glycol',
            'Various hydrocarbons'
        ],
        'Purpose': [
            'Mercaptan removal',
            'Acid gas removal',
            'Water removal',
            'Component separation'
        ]
    }
    
    # Safety data sheet
    safety_data = {
        'Chemical': [
            'Monoethanolamine', 'Diethanolamine', 'Triethanolamine',
            'Hydrogen Sulfide', 'Mercaptans'
        ],
        'Hazard Class': ['Corrosive', 'Corrosive', 'Irritant', 'Toxic Gas', 'Toxic'],
        'Flash Point (°C)': [93, 137, 179, -82, -15],
        'PPE Required': [
            'Gloves, goggles, apron',
            'Gloves, goggles, apron', 
            'Gloves, goggles',
            'SCSR, gas detector',
            'Gas detector, ventilation'
        ],
        'First Aid': [
            'Flush with water, seek medical attention',
            'Flush with water, seek medical attention',
            'Flush with water',
            'Move to fresh air, oxygen if needed',
            'Move to fresh air, wash contaminated areas'
        ]
    }
    
    # Create Excel file with multiple sheets
    with pd.ExcelWriter('sample_chemical_database.xlsx', engine='openpyxl') as writer:
        pd.DataFrame(chemicals_data).to_excel(writer, sheet_name='Chemical Properties', index=False)
        pd.DataFrame(process_data).to_excel(writer, sheet_name='Process Conditions', index=False)
        pd.DataFrame(safety_data).to_excel(writer, sheet_name='Safety Data', index=False)
    
    print("✅ Created sample_chemical_database.xlsx")

def create_sample_powerpoint():
    """Create a sample PowerPoint presentation"""
    prs = Presentation()
    
    # Slide 1: Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Chemical Process Training"
    subtitle.text = "Amine Treatment and Gas Processing"
    
    # Slide 2: Amine Chemistry
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = 'Amine Chemistry Overview'
    
    tf = body_shape.text_frame
    tf.text = 'Primary Amines'
    
    p = tf.add_paragraph()
    p.text = 'Monoethanolamine (MEA) - Most reactive, highest heat of reaction'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Used for deep acid gas removal'
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = 'Secondary Amines'
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = 'Diethanolamine (DEA) - Lower heat of reaction than MEA'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Good for moderate acid gas concentrations'
    p.level = 2
    
    # Slide 3: Process Safety
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = 'Ethanolamine Safety Considerations'
    
    tf = body_shape.text_frame
    tf.text = 'Health Hazards'
    
    p = tf.add_paragraph()
    p.text = 'Corrosive to skin and eyes'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Can cause respiratory irritation'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Safety Measures'
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = 'Always wear appropriate PPE'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Ensure adequate ventilation'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Have emergency eyewash stations nearby'
    p.level = 1
    
    # Slide 4: Process Flow
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = 'Amine Treatment Process Flow'
    
    tf = body_shape.text_frame
    tf.text = 'Gas feed enters absorber column'
    
    p = tf.add_paragraph()
    p.text = 'Contact with lean amine solution'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Acid gases absorbed into amine'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Rich amine sent to regenerator'
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = 'Heat applied to strip acid gases'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Lean amine recycled back to absorber'
    p.level = 1
    
    prs.save('sample_process_training.pptx')
    print("✅ Created sample_process_training.pptx")

def create_sample_text_file():
    """Create a sample text file"""
    content = """
Chemical Process Equipment Manual
==================================

Chapter 5: Amine Treatment Systems

Introduction:
Amine treatment is a widely used method for removing acid gases (H2S and CO2) from natural gas and refinery gas streams. The process utilizes aqueous solutions of various amines, with ethanolamine being one of the most common.

Types of Ethanolamines:

1. Monoethanolamine (MEA)
   - Chemical formula: HOCH2CH2NH2
   - Most reactive amine
   - Highest heat of reaction
   - Used for very deep acid gas removal
   - Requires more energy for regeneration

2. Diethanolamine (DEA)
   - Chemical formula: (HOCH2CH2)2NH
   - Less reactive than MEA
   - Lower heat of reaction
   - Better thermal stability
   - Suitable for moderate acid gas concentrations

3. Methyldiethanolamine (MDEA)
   - Chemical formula: (HOCH2CH2)2NCH3
   - Selective for H2S over CO2
   - Lowest heat of reaction
   - Excellent thermal stability
   - Used in selective treating applications

Process Description:
The amine treatment process consists of two main sections:
1. Absorption section - where acid gases are absorbed
2. Regeneration section - where the amine is stripped and recycled

Equipment Specifications:
- Absorber tower: typically 20-40 trays
- Regenerator: 18-24 trays
- Heat exchanger: shell and tube design
- Pumps: centrifugal type with corrosion-resistant materials

Safety Considerations:
When working with ethanolamine solutions:
- Always wear appropriate PPE including gloves and safety glasses
- Ensure adequate ventilation in work areas
- Have emergency shower and eyewash stations readily available
- Monitor for amine vapors in confined spaces
- Follow proper lockout/tagout procedures during maintenance

Troubleshooting:
Common issues in amine systems include:
- Foaming due to contamination
- Corrosion from acid gas loading
- Thermal degradation at high temperatures
- Mechanical problems with pumps and heat exchangers

For more information, consult the process flow diagrams and P&ID drawings in the appendix.
"""
    
    with open('sample_equipment_manual.txt', 'w') as f:
        f.write(content)
    
    print("✅ Created sample_equipment_manual.txt")

def main():
    """Create all sample documents"""
    print("🔧 Creating sample documents for testing...")
    print("-" * 50)
    
    try:
        create_sample_word_doc()
        create_sample_excel_file() 
        create_sample_powerpoint()
        create_sample_text_file()
        
        print("-" * 50)
        print("✅ All sample documents created successfully!")
        print("\nSample files created:")
        print("📄 sample_chemical_manual.docx")
        print("📊 sample_chemical_database.xlsx") 
        print("📈 sample_process_training.pptx")
        print("📝 sample_equipment_manual.txt")
        print("\nYou can now test the multi-format RAG system with these files.")
        
    except Exception as e:
        print(f"❌ Error creating sample documents: {e}")

if __name__ == "__main__":
    main()