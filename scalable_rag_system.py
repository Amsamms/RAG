#!/usr/bin/env python3
"""
Scalable RAG System for Multiple Document Formats
Supports: PDF, Word, Excel, PowerPoint + Hundreds of documents
"""

import os
import fitz  # PyMuPDF for PDFs
import pandas as pd  # For Excel
import docx  # For Word documents
from pptx import Presentation  # For PowerPoint
import chromadb
from sentence_transformers import SentenceTransformer
from typing import List, Dict, Any, Optional
from pathlib import Path
import openai
from dotenv import load_dotenv
import hashlib
import json
from tqdm import tqdm
import logging
from concurrent.futures import ThreadPoolExecutor, as_completed
import time

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Load environment variables
load_dotenv()

class ScalableMultiFormatRAG:
    def __init__(self, collection_name: str = "multi_format_documents", batch_size: int = 50):
        """Initialize scalable RAG system for multiple document formats"""
        self.client = chromadb.Client()
        self.collection_name = collection_name
        self.batch_size = batch_size
        self.embedding_model_name = os.getenv("EMBEDDING_MODEL", "all-MiniLM-L6-v2")
        self.embedding_model = SentenceTransformer(self.embedding_model_name)
        
        # OpenAI configuration
        self.openai_api_key = os.getenv("OPENAI_API_KEY")
        self.openai_model = os.getenv("OPENAI_MODEL", "gpt-3.5-turbo")
        
        if self.openai_api_key:
            openai.api_key = self.openai_api_key
            self.llm_available = True
        else:
            self.llm_available = False
            logger.warning("OpenAI API key not found. LLM features disabled.")
        
        # Initialize collection
        try:
            self.collection = self.client.get_collection(name=collection_name)
            logger.info(f"Using existing collection: {collection_name}")
        except Exception:
            self.collection = self.client.create_collection(name=collection_name)
            logger.info(f"Created new collection: {collection_name}")
        
        # Document processing statistics
        self.processing_stats = {
            'total_files': 0,
            'successful': 0,
            'failed': 0,
            'total_chunks': 0,
            'processing_time': 0
        }
    
    def extract_text_from_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract text from PDF files"""
        try:
            doc = fitz.open(file_path)
            pages_data = []
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                
                if text.strip():
                    chunks = self._split_text(text)
                    for chunk_id, chunk in enumerate(chunks):
                        if chunk.strip():
                            pages_data.append({
                                'text': chunk.strip(),
                                'document': os.path.basename(file_path),
                                'page': page_num + 1,
                                'chunk_id': chunk_id,
                                'file_path': file_path,
                                'file_type': 'pdf'
                            })
            
            doc.close()
            return pages_data
            
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {e}")
            return []
    
    def extract_text_from_word(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract text from Word documents"""
        try:
            doc = docx.Document(file_path)
            full_text = []
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text)
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells])
                    if row_text.strip():
                        full_text.append(row_text)
            
            # Combine and chunk
            combined_text = "\n".join(full_text)
            chunks = self._split_text(combined_text)
            
            pages_data = []
            for chunk_id, chunk in enumerate(chunks):
                if chunk.strip():
                    pages_data.append({
                        'text': chunk.strip(),
                        'document': os.path.basename(file_path),
                        'page': 1,  # Word documents don't have pages like PDFs
                        'chunk_id': chunk_id,
                        'file_path': file_path,
                        'file_type': 'word'
                    })
            
            return pages_data
            
        except Exception as e:
            logger.error(f"Error processing Word document {file_path}: {e}")
            return []
    
    def extract_text_from_excel(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract text from Excel files"""
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            pages_data = []
            
            for sheet_num, sheet_name in enumerate(excel_file.sheet_names):
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # Convert DataFrame to text
                sheet_text = f"Sheet: {sheet_name}\n"
                sheet_text += df.to_string(index=False)
                
                chunks = self._split_text(sheet_text)
                for chunk_id, chunk in enumerate(chunks):
                    if chunk.strip():
                        pages_data.append({
                            'text': chunk.strip(),
                            'document': os.path.basename(file_path),
                            'page': sheet_num + 1,  # Use sheet number as page
                            'chunk_id': chunk_id,
                            'file_path': file_path,
                            'file_type': 'excel',
                            'sheet_name': sheet_name
                        })
            
            return pages_data
            
        except Exception as e:
            logger.error(f"Error processing Excel file {file_path}: {e}")
            return []
    
    def extract_text_from_powerpoint(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract text from PowerPoint files"""
        try:
            prs = Presentation(file_path)
            pages_data = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = f"Slide {slide_num + 1}:\n"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                
                chunks = self._split_text(slide_text)
                for chunk_id, chunk in enumerate(chunks):
                    if chunk.strip():
                        pages_data.append({
                            'text': chunk.strip(),
                            'document': os.path.basename(file_path),
                            'page': slide_num + 1,
                            'chunk_id': chunk_id,
                            'file_path': file_path,
                            'file_type': 'powerpoint'
                        })
            
            return pages_data
            
        except Exception as e:
            logger.error(f"Error processing PowerPoint file {file_path}: {e}")
            return []
    
    def _split_text(self, text: str, chunk_size: int = 500) -> List[str]:
        """Split text into chunks optimized for large-scale processing"""
        paragraphs = text.split('\n\n')
        chunks = []
        
        for paragraph in paragraphs:
            if len(paragraph) <= chunk_size:
                chunks.append(paragraph)
            else:
                # Split long paragraphs by sentences
                sentences = paragraph.split('. ')
                current_chunk = ""
                
                for sentence in sentences:
                    if len(current_chunk + sentence) <= chunk_size:
                        current_chunk += sentence + ". "
                    else:
                        if current_chunk:
                            chunks.append(current_chunk.strip())
                        current_chunk = sentence + ". "
                
                if current_chunk:
                    chunks.append(current_chunk.strip())
        
        return chunks
    
    def process_single_file(self, file_path: Path) -> List[Dict[str, Any]]:
        """Process a single file based on its extension"""
        file_ext = file_path.suffix.lower()
        
        if file_ext == '.pdf':
            return self.extract_text_from_pdf(str(file_path))
        elif file_ext in ['.docx', '.doc']:
            return self.extract_text_from_word(str(file_path))
        elif file_ext in ['.xlsx', '.xls']:
            return self.extract_text_from_excel(str(file_path))
        elif file_ext in ['.pptx', '.ppt']:
            return self.extract_text_from_powerpoint(str(file_path))
        else:
            logger.warning(f"Unsupported file type: {file_ext}")
            return []
    
    def process_documents_parallel(self, directory: str = ".", max_workers: int = 4) -> int:
        """Process documents in parallel for better performance"""
        start_time = time.time()
        
        # Find all supported files
        supported_extensions = {'.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt'}
        all_files = []
        
        for ext in supported_extensions:
            all_files.extend(Path(directory).glob(f"*{ext}"))
        
        self.processing_stats['total_files'] = len(all_files)
        logger.info(f"Found {len(all_files)} files to process")
        
        if not all_files:
            logger.warning("No supported files found")
            return 0
        
        # Process files in parallel
        all_documents = []
        all_embeddings = []
        all_metadatas = []
        all_ids = []
        
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            # Submit all files for processing
            future_to_file = {executor.submit(self.process_single_file, file_path): file_path 
                             for file_path in all_files}
            
            # Collect results with progress bar
            for future in tqdm(as_completed(future_to_file), total=len(all_files), desc="Processing files"):
                file_path = future_to_file[future]
                try:
                    pages_data = future.result()
                    if pages_data:
                        self.processing_stats['successful'] += 1
                        
                        # Process chunks in batches for memory efficiency
                        for page_data in pages_data:
                            # Create embeddings
                            embedding = self.embedding_model.encode(page_data['text'])
                            
                            # Prepare data for ChromaDB
                            doc_id = f"{page_data['document']}_page_{page_data['page']}_chunk_{page_data['chunk_id']}"
                            
                            all_documents.append(page_data['text'])
                            all_embeddings.append(embedding.tolist())
                            all_metadatas.append({
                                'document': page_data['document'],
                                'page': page_data['page'],
                                'chunk_id': page_data['chunk_id'],
                                'file_path': page_data['file_path'],
                                'file_type': page_data['file_type']
                            })
                            all_ids.append(doc_id)
                            
                            # Add to database in batches
                            if len(all_documents) >= self.batch_size:
                                self._add_batch_to_db(all_documents, all_embeddings, all_metadatas, all_ids)
                                all_documents, all_embeddings, all_metadatas, all_ids = [], [], [], []
                    
                except Exception as e:
                    logger.error(f"Error processing {file_path}: {e}")
                    self.processing_stats['failed'] += 1
        
        # Add remaining documents
        if all_documents:
            self._add_batch_to_db(all_documents, all_embeddings, all_metadatas, all_ids)
        
        # Update statistics
        self.processing_stats['total_chunks'] = self.collection.count()
        self.processing_stats['processing_time'] = time.time() - start_time
        
        logger.info(f"Processing complete: {self.processing_stats}")
        return self.processing_stats['total_chunks']
    
    def _add_batch_to_db(self, documents, embeddings, metadatas, ids):
        """Add a batch of documents to the database"""
        try:
            self.collection.add(
                documents=documents,
                embeddings=embeddings,
                metadatas=metadatas,
                ids=ids
            )
        except Exception as e:
            logger.error(f"Error adding batch to database: {e}")
    
    def search_documents(self, question: str, n_results: int = 10) -> Dict[str, Any]:
        """Search documents with optimized performance for large datasets"""
        question_embedding = self.embedding_model.encode(question)
        
        results = self.collection.query(
            query_embeddings=[question_embedding.tolist()],
            n_results=n_results
        )
        
        response = {
            'question': question,
            'results': []
        }
        
        if results['documents'] and results['documents'][0]:
            for i in range(len(results['documents'][0])):
                result_item = {
                    'text': results['documents'][0][i],
                    'document': results['metadatas'][0][i]['document'],
                    'page': results['metadatas'][0][i]['page'],
                    'file_type': results['metadatas'][0][i]['file_type'],
                    'distance': results['distances'][0][i] if 'distances' in results else None
                }
                response['results'].append(result_item)
        
        return response
    
    def get_cost_estimate(self, query_count: int = 100) -> Dict[str, Any]:
        """Estimate OpenAI costs for given usage"""
        
        # Current OpenAI pricing (as of 2024)
        pricing = {
            'gpt-3.5-turbo': {
                'input': 0.0015,   # per 1K tokens
                'output': 0.002    # per 1K tokens
            },
            'gpt-4': {
                'input': 0.03,     # per 1K tokens
                'output': 0.06     # per 1K tokens
            }
        }
        
        # Estimate tokens per query
        avg_context_tokens = 2000  # Average context from search results
        avg_query_tokens = 20      # Average question length
        avg_response_tokens = 200  # Average response length
        
        total_input_tokens = (avg_context_tokens + avg_query_tokens) * query_count
        total_output_tokens = avg_response_tokens * query_count
        
        cost_estimate = {}
        for model, prices in pricing.items():
            input_cost = (total_input_tokens / 1000) * prices['input']
            output_cost = (total_output_tokens / 1000) * prices['output']
            total_cost = input_cost + output_cost
            
            cost_estimate[model] = {
                'input_cost': round(input_cost, 4),
                'output_cost': round(output_cost, 4),
                'total_cost': round(total_cost, 4),
                'cost_per_query': round(total_cost / query_count, 6)
            }
        
        return cost_estimate

# Installation requirements for additional formats
def install_requirements():
    """Install additional packages for multi-format support"""
    requirements = [
        "python-docx",  # For Word documents
        "python-pptx",  # For PowerPoint
        "pandas",       # For Excel (already installed)
        "openpyxl"      # For Excel support
    ]
    
    print("📦 Additional packages needed for multi-format support:")
    for req in requirements:
        print(f"   pip install {req}")

if __name__ == "__main__":
    # Test the scalable system
    rag = ScalableMultiFormatRAG()
    
    # Show cost estimates
    costs = rag.get_cost_estimate(100)
    print("💰 Cost Estimates for 100 queries:")
    for model, cost_info in costs.items():
        print(f"   {model}: ${cost_info['total_cost']} (${cost_info['cost_per_query']}/query)")
    
    install_requirements()